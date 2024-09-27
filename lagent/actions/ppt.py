from typing import Dict, Optional, Type

from aioify import aioify

from lagent.actions.base_action import AsyncActionMixin, BaseAction, tool_api
from lagent.actions.parser import BaseParser, JsonParser

THEME_MAPPING = {
    'Default': {
        'template': None,
        'title': 'Title Slide',
        'single': 'Title and Content',
        'two': 'Two Content',
    }
}


class PPT(BaseAction):
    """Plugin to create ppt slides with text, paragraph, images in good looking styles."""

    def __init__(
        self,
        theme_mapping: Optional[Dict[str, dict]] = None,
        description: Optional[dict] = None,
        parser: Type[BaseParser] = JsonParser,
    ):
        super().__init__(description, parser)
        self.theme_mapping = theme_mapping or THEME_MAPPING
        self.pointer = None
        self.location = None

    @tool_api(explode_return=True)
    def create_file(self, theme: str, abs_location: str) -> dict:
        """Create a pptx file with specific themes.

        Args:
            theme (:class:`str`): the theme used. The value should be one of ['Default'].
            abs_location (:class:`str`): the ppt file's absolute location

        Returns:
            :class:`dict`: operation status
                * status: the result of the execution
        """
        from pptx import Presentation
        self.location = abs_location
        try:
            self.pointer = Presentation(self.theme_mapping[theme]['template'])
            self.pointer.slide_master.name = theme
            # print('created')
        except Exception as e:
            print(e)
        return dict(status='created a ppt file.')

    @tool_api(explode_return=True)
    def add_first_page(self, title: str, subtitle: str) -> dict:
        """Add the first page of ppt.

        Args:
            title (:class:`str`): the title of ppt
            subtitle (:class:`str`): the subtitle of ppt

        Returns:
            :class:`dict`: operation status
                * status: the result of the execution
        """
        layout_name = self.theme_mapping[
            self.pointer.slide_master.name]['title']
        layout = next(i for i in self.pointer.slide_master.slide_layouts
                      if i.name == layout_name)
        slide = self.pointer.slides.add_slide(layout)
        ph_title, ph_subtitle = slide.placeholders
        ph_title.text = title
        if subtitle:
            ph_subtitle.text = subtitle
        return dict(status='added page')

    @tool_api(explode_return=True)
    def add_text_page(self, title: str, bullet_items: str) -> dict:
        """Add text page of ppt.

        Args:
            title (:class:`str`): the title of the page
            bullet_items (:class:`str`): bullet_items should be string, for multiple bullet items, please use [SPAN] to separate them.

        Returns:
            :class:`dict`: operation status
                * status: the result of the execution
        """
        layout_name = self.theme_mapping[
            self.pointer.slide_master.name]['single']
        layout = next(i for i in self.pointer.slide_master.slide_layouts
                      if i.name == layout_name)
        slide = self.pointer.slides.add_slide(layout)
        ph_title, ph_body = slide.placeholders
        ph_title.text = title
        ph = ph_body
        tf = ph.text_frame
        for i, item in enumerate(bullet_items.split('[SPAN]')):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item.strip()
            p.level = 0
        return dict(status='added page')

    @tool_api(explode_return=True)
    def add_text_image_page(self, title: str, bullet_items: str,
                            image: str) -> dict:
        """Add a text page with one image. Image should be a path.

        Args:
            title (:class:`str`): the title of the page
            bullet_items (:class:`str`): bullet_items should be string, for multiple bullet items, please use [SPAN] to separate them.
            image (:class:`str`): the path of the image

        Returns:
            :class:`dict`: operation status
                * status: the result of the execution
        """
        from PIL import Image
        layout_name = self.theme_mapping[self.pointer.slide_master.name]['two']
        layout = next(i for i in self.pointer.slide_master.slide_layouts
                      if i.name == layout_name)
        slide = self.pointer.slides.add_slide(layout)
        ph_title, ph_body1, ph_body2 = slide.placeholders
        ph_title.text = title
        ph = ph_body2
        image = Image.open(image)
        image_pil = image.to_pil()
        left = ph.left
        width = ph.width
        height = int(width / image_pil.width * image_pil.height)
        top = (ph.top + (ph.top + ph.height)) // 2 - height // 2
        slide.shapes.add_picture(image.to_path(), left, top, width, height)

        ph = ph_body1
        tf = ph.text_frame
        for i, item in enumerate(bullet_items.split('[SPAN]')):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item.strip()
            p.level = 0

        return dict(status='added page')

    @tool_api(explode_return=True)
    def submit_file(self) -> dict:
        """When all steps done, YOU MUST use submit_file() to submit your work.

        Returns:
            :class:`dict`: operation status
                * status: the result of the execution
        """
        # file_path = os.path.join(self.CACHE_DIR, f'{self._return_timestamp()}.pptx')
        # self.pointer.save(file_path)
        # retreival_url = upload_file(file_path)
        self.pointer.save(self.location)
        return dict(status=f'submitted. view ppt at {self.location}')


class AsyncPPT(AsyncActionMixin, PPT):
    """Plugin to create ppt slides with text, paragraph, images in good looking styles."""

    @tool_api(explode_return=True)
    @aioify
    def create_file(self, theme: str, abs_location: str) -> dict:
        """Create a pptx file with specific themes.

        Args:
            theme (:class:`str`): the theme used. The value should be one of ['Default'].
            abs_location (:class:`str`): the ppt file's absolute location

        Returns:
            :class:`dict`: operation status
                * status: the result of the execution
        """
        return super().create_file(theme, abs_location)

    @tool_api(explode_return=True)
    @aioify
    def add_first_page(self, title: str, subtitle: str) -> dict:
        """Add the first page of ppt.

        Args:
            title (:class:`str`): the title of ppt
            subtitle (:class:`str`): the subtitle of ppt

        Returns:
            :class:`dict`: operation status
                * status: the result of the execution
        """
        return super().add_first_page(title, subtitle)

    @tool_api(explode_return=True)
    @aioify
    def add_text_page(self, title: str, bullet_items: str) -> dict:
        """Add text page of ppt.

        Args:
            title (:class:`str`): the title of the page
            bullet_items (:class:`str`): bullet_items should be string, for multiple bullet items, please use [SPAN] to separate them.

        Returns:
            :class:`dict`: operation status
                * status: the result of the execution
        """
        return super().add_text_page(title, bullet_items)

    @tool_api(explode_return=True)
    @aioify
    def add_text_image_page(self, title: str, bullet_items: str,
                            image: str) -> dict:
        """Add a text page with one image. Image should be a path.

        Args:
            title (:class:`str`): the title of the page
            bullet_items (:class:`str`): bullet_items should be string, for multiple bullet items, please use [SPAN] to separate them.
            image (:class:`str`): the path of the image

        Returns:
            :class:`dict`: operation status
                * status: the result of the execution
        """
        return super().add_text_image_page(title, bullet_items, image)

    @tool_api(explode_return=True)
    @aioify
    def submit_file(self) -> dict:
        """When all steps done, YOU MUST use submit_file() to submit your work.

        Returns:
            :class:`dict`: operation status
                * status: the result of the execution
        """
        return super().submit_file()
